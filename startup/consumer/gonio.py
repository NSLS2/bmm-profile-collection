
import numpy, os, sys, pandas, pathlib, datetime, re
from bluesky import __version__ as bluesky_version

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import datetime

from tools import echo_slack, experiment_folder, file_resource, profile_configuration

def log_entry(logger, message):
    #if logger.name == 'BMM file manager logger' or logger.name == 'bluesky_kafka':
    #print(message)
    #post_to_slack(message)
    #echo_slack(text = message,
    #           icon = 'message',
    #           rid  = None )
    logger.info(message)



class XRRFile():


    # binary                0  1        2                  3                                         4
    # level                 0  1        2        3         4         5         6         7*          8
    measured_attenuation = [1, 6.85865, 47.0088, 318.6107, 2225.346, 15046.19, 97500.05, 668718.718, ]

    
    def to_xdi(self, catalog=None, uid=None, stub=None, logger=None):
        '''Write an XDI-style file for an XRR scan.

        '''
        metadata = catalog[uid].metadata
        xdi = metadata["start"]["XDI"]
        if stub is None:
            stub = xdi['_filename']
        fname = os.path.join(experiment_folder(catalog, uid), stub+'.xdi')
        handle = open(fname, 'w')
        handle.write(f'# XDI/1.0 BlueSky/{bluesky_version} BMM/{pathlib.Path(sys.executable).parts[-3]}\n')
        
        ## header lines with metadata from the XDi dictionary
        for family in ('Beamline', 'Detector', 'Element', 'Facility', 'Mono', 'Sample', 'Scan'):
            for k in xdi[family].keys():
                if family == 'Sample' and k == 'comment':
                    continue
                if family == 'Sample' and k == 'extra_metadata':
                    continue
                handle.write(f'# {family}.{k}: {xdi[family][k]}\n')
        start = datetime.datetime.fromtimestamp(metadata['start']['time']).strftime("%Y-%m-%dT%H:%M:%S") # '%A, %d %B, %Y %I:%M %p')
        end   = datetime.datetime.fromtimestamp(metadata['stop']['time']).strftime("%Y-%m-%dT%H:%M:%S") # '%A, %d %B, %Y %I:%M %p')
        handle.write(f'# Scan.start_time: {start}\n')
        handle.write(f'# Scan.end_time: {end}\n')
        handle.write(f'# Scan.uid: {uid}\n')
        handle.write(f'# Scan.transient_id: {metadata["start"]["scan_id"]}\n')

        if 'mythen-2' in metadata['start']['detectors']:
            hdf5files = file_resource(catalog, uid)
            for h in hdf5files:
                relative = '/'.join(h.split('/')[-6:])
                if 'mythen' in relative:
                    handle.write(f'# Scan.mythen_hdf5_file: {relative}\n')

        #handle.write( '# Scan.plot_hint: \n')
        handle.write( '# Column.1: eta degrees\n')
        handle.write( '# Column.2: delta degrees\n')
        handle.write( '# Column.3: measurement_time seconds\n')
        handle.write( '# Column.4: monitor counts\n')
        handle.write( '# Column.5: mca_full counts\n')
        handle.write( '# Column.6: mca_narrow counts\n')
        handle.write( '# Column.7: attenuator\n')

        ## Column.N header lines
        column_list = ['eta', 'delta', 'dwti_dwell_time', 'monitor', 'mca_full', 'mca_narrow', 'attenuator_attenuation']
        column_labels = ['eta', 'delta', 'measurement_time', 'monitor', 'mca_full', 'mac_narrow', 'attenuator']

        xa = catalog[uid].primary.read(column_list)
        p = xa.to_pandas()
        
        ## use eta as the pandas index
        p.set_index('eta')

        ## comment and separator lines
        handle.write('# //////////////////////////////////////////////////////////\n')
        if '_comment' in xdi:
            for l in xdi["_comment"]:
                handle.write(f'# {l}\n')
            else:
                handle.write(f'# \n')
            handle.write('# ----------------------------------------------------------\n')
        handle.write('# ')

        ## dump the data table and close the file
        handle.write(p.to_csv(None, sep=' ', columns=column_list, index=False, header=column_labels, float_format='%.6f'))
        handle.flush()
        handle.close()

        log_entry(logger, f'wrote XRR data to {fname}')


    def to_txt(self, catalog=None, uid=None, stub=None, logger=None, style='short'):

        header = '''% Description of the file: <name of detector> <columns>
%                      or: <name of detector> <first column> <last column>
%
%D	Delta	1
%D	Eta	2
%D	Nu	3
%D	Wheel 1	4
%D	mca	5
%D	dir	6
%D	Monitor	7
%D	Seconds	8
'''
        longheader = '%D	LinearDetector	9	1288\n'

        nuval = catalog[uid].baseline['data']['nu'][0]
        column_list = ['delta', 'eta', 'attenuator_attenuation', 'mca_full', 'mca_narrow', 'monitor', 'dwti_dwell_time']
        xa = catalog[uid].primary.read(column_list)
        p = xa.to_pandas()
        column_list.insert(2, 'nu')
        npoints = len(catalog[uid].primary['data']['eta'])
        nu = nuval * numpy.ones(npoints)
        p['nu'] = nu


        if style in ('short', 'both'):
            fname = os.path.join(experiment_folder(catalog, uid), stub+'_short.txt')
            handle = open(fname, 'w')
            handle.write(header)
            handle.write('\n')
            handle.write(p.to_csv(None, sep=' ', columns=column_list, index=False, header=False, float_format='%.6f'))
            handle.flush()
            handle.close()
            
            log_entry(logger, f'wrote XRR data to {fname}')

        
        if style in ('long', 'both'):
            fullmca = catalog[uid].primary['data']['mythen-2_image'][:,0,:].astype(int)
            mcabins = list((f'bin{i+1}' for i in range(fullmca.shape[-1]) ))
            mcaFrame = pandas.DataFrame(fullmca, columns=mcabins)
            
            fname = os.path.join(experiment_folder(catalog, uid), stub+'_long.txt')
            handle = open(fname, 'w')
            handle.write(header)
            handle.write(longheader)
            handle.write('\n')

            p = p.join(mcaFrame)
            handle.write(p.to_csv(None, sep=' ', columns=column_list+mcabins, index=False, header=False, float_format='%.6f'))
            handle.flush()
            handle.close()

            log_entry(logger, f'wrote XRR data to {fname}')


    def unclobbered_filename(self, filename):
        if os.path.exists(filename) is False:
            return filename
        name, extension = os.path.splitext(filename)
        pattern = re.compile('\((\d+)\)$')
        s = pattern.search(name)
        if s is None:
            name = name + '(1)'
        else:
            was = s.group()
            next_index = int(s.groups()[0]) + 1
            name = name.replace(was, f'({next_index})')
        return name+extension

        
    def linescan_file(self, catalog=None, uid=None, stub=None, motor=None, detector=None, logger=None):
        header = f'''% Description of the file: <name of detector> <columns>
%                      or: <name of detector> <first column> <last column>
%
%D	{motor.capitalize()}	1
'''
        n = 2
        if detector.lower() == 'mythen':
            header += '%D	mca_full	2\n'
            header += '%D	mca_narrow	3\n'
            column_list = [motor, 'mca_full', 'mca_narrow', 'monitor', 'dwti_dwell_time']
            n = 3
        elif detector.lower() == 'mca_full':
            header += '%D	mca_full	2\n'
            column_list = [motor, 'mca_full', 'monitor', 'dwti_dwell_time']
            n = 2
        elif detector.lower() == 'mca_narrow':
            header += '%D	mca_narrow	2\n'
            column_list = [motor, 'mca_narrow', 'monitor', 'dwti_dwell_time']
            n = 2

        header += f'%D	Monitor	{n+1}\n'
        header += f'%D	Seconds	{n+2}\n'

        header += f'%D	LinearDetector	{n+3}	{n+1283}\n'

        xa = catalog[uid].primary.read(column_list)
        p = xa.to_pandas()
        
        fullmca = catalog[uid].primary['data']['mythen-2_image'][:,0,:].astype(int)
        mcabins = list((f'bin{i+1}' for i in range(fullmca.shape[-1]) ))
        mcaFrame = pandas.DataFrame(fullmca, columns=mcabins)

        fname = os.path.join(experiment_folder(catalog, uid), stub+'.dat')
        fname = self.unclobbered_filename(fname)
        handle = open(fname, 'w')
        handle.write(header)

        p = p.join(mcaFrame)
        handle.write(p.to_csv(None, sep=' ', columns=column_list+mcabins, index=False, header=False, float_format='%.6f'))
        handle.flush()
        handle.close()

        log_entry(logger, f'wrote XRR linescan to {fname}')
            
    def calibration_file(self, catalog=None, uid=None, stub=None, motor=None, detector=None, logger=None):
        header = f'''% Description of the file: <name of detector> <columns>
%                      or: <name of detector> <first column> <last column>
%
%D	Delta	1
%D	Eta	2
%D	Monitor	3
%D	LinearDetector  4	1283

'''
        etaval = catalog[uid].baseline['data']['eta'][0]
        column_list = ['delta', 'monitor']
        xa = catalog[uid].primary.read(column_list)
        p = xa.to_pandas()
        column_list.insert(2, 'eta')
        npoints = len(catalog[uid].primary['data']['delta'])
        eta = etaval * numpy.ones(npoints)
        p['eta'] = eta

        fullmca = catalog[uid].primary['data']['mythen-2_image'][:,0,:].astype(int)
        mcabins = list((f'bin{i+1}' for i in range(fullmca.shape[-1]) ))
        mcaFrame = pandas.DataFrame(fullmca, columns=mcabins)

        

        fname = os.path.join(experiment_folder(catalog, uid), stub+'.dat')
        fname = self.unclobbered_filename(fname)
        handle = open(fname, 'w')
        handle.write(header)

        p = p.join(mcaFrame)
        handle.write(p.to_csv(None, sep=' ', columns=column_list+mcabins, index=False, header=False, float_format='%.6f'))
        handle.flush()
        handle.close()

        log_entry(logger, f'wrote XRR data to {fname}')


        
    def mythen_calibration(self, catalog=None, uid=None, path=None, now=None, stamp=None, setup=None, gap=None,
                           energy=8600, pixel0=None, angle_per_pixel=None, stub=None, dw=0.12,
                           rw=0.2, slits_b=0.3, slits_i=0.5, slits_o=0.5, slits_t=0.3, logger=None):  # fixme! fitA, fitB, fitC
        '''Write a PowerPoint summary of the calibration using the established
        layout of the report in use by the IBM folks.

        '''

        ## fixme!
        fitA = fitB = fitC = 0

        ## make a pptx with a single blank slide and no placeholders
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)


        ## make a Text box at the top with the path to the proposal
        ## folder and the UID of the calibration scan
        top = Inches(0.2)
        left = Inches(1.5)
        width = Inches(1)
        height = Inches(1)
        txBox1 = slide.shapes.add_textbox(left, top, width, height)
        tf1 = txBox1.text_frame
        p = tf1.paragraphs[0]
        run = p.add_run()
        run.text = f'{path}\n{uid}'
        run.font.size = Pt(10)

        ## make a Text box for the all the header information, date,
        ## measurement type, gap, energy, calibration fit result,
        ## center pixel position, detector distance calculation
        top = Inches(0.75)
        left = Inches(4)
        width = Inches(1)
        height = Inches(1)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame

        #p = tf.add_paragraph()
        tf2.text = f"{now} calibration"

        p = tf2.add_paragraph()
        p.text = f'(for {setup})'

        p = tf2.add_paragraph()
        p.text = f'Gap={gap} mm, E= {energy} keV'

        p = tf2.add_paragraph()
        p.text = f'FIT: arctan((channel - {pixel0})/{angle_per_pixel})'

        p = tf2.add_paragraph()
        p.text = f'PIXEL 0 = {pixel0}; D=0.05 x {angle_per_pixel} = {angle_per_pixel/0.05} mm'

        ## justify the first two text boxes
        for para in tf1.paragraphs:
            para.alignment = PP_ALIGN.LEFT
        for para in tf2.paragraphs:
            para.alignment = PP_ALIGN.CENTER


        ## make a box for the picture of the fit
        left=Inches(4)
        top=Inches(2.5)
        width=Inches(6)
        height=Inches(4)
        pic = slide.shapes.add_picture(os.path.join(path, 'snapshots', stub+'.png'),
                                       left, top, width=width, height=height)

        
        ## make a box for the mythen ROI settings
        left=Inches(0.9)
        top=Inches(3)
        width=Inches(3)
        height=Inches(2)

        ROIBox = slide.shapes.add_textbox(left, top, width, height)
        roitext = ROIBox.text_frame
        roitext.clear()
        p = roitext.paragraphs[0]
        run = p.add_run()
        run.text = '''ROIs:
        dir = +/-{dw}, pixels {pixel0-dw}-{pixel0+dw}, {0.05*(2*dw+1)}mm
        refl = +/-{rw}, pixels {pixel0-rw}-{pixel0+rw}, {0.05*(2*rw+1)}mm
        '''
        run.font.size = Pt(10)


        ## make a box for the slit settings
        left=Inches(0.6)
        top=Inches(4.5)
        width=Inches(3)
        height=Inches(2)

        SlitsBox = slide.shapes.add_textbox(left, top, width, height)
        slitstext = SlitsBox.text_frame
        slitstext.clear()
        p = slitstext.paragraphs[0]
        run = p.add_run()
        run.text = '''Incident slits:
        s1t, s1b = {slits_biot[0]}, V={2*slits_biot[0]}
        s1o, s1i = {slits_biot[1]}, H={2*slits_biot[1]}
        '''
        run.font.size = Pt(14)



        if 'pole' in setup.lower():
            left=Inches(3)
            top=Inches(6.5)
            width=Inches(6)
            height=Inches(1)

            CalBox = slide.shapes.add_textbox(left, top, width, height)
            caltext = CalBox.text_frame
            caltext.clear()
            p = caltext.paragraphs[0]
            run = p.add_run()
            run.text = 'CHESS calibration: 2θ = ({fitA}*pixel² + {fitB}*pixel + {fitC}) + del'
            run.font.size = Pt(14)


        run.font.size = Pt(10)

        fname = os.path.join(path, f'{stub}_{stamp}.pptx')
        prs.save(fname)
        log_entry(logger, f'wrote Mythen calibration report to {fname}')
